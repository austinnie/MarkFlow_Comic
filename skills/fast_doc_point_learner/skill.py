"""
fast_doc_point_learner - 文档快速学习助手

功能：
  - 多格式支持 (Word, PDF, PPT, Excel, TXT)
  - 内容提取 (标题、段落、表格、图片说明)
  - 智能摘要 (使用 AI 生成)
  - 关系分析 (文档间引用、依赖、补充关系)
  - 结构化展示 (文档内容和关系图谱)
  - 多语言支持 (zh/en/ja/ko/es/fr/de)
"""

import os
import re
import json
import logging
from pathlib import Path
from typing import Dict, Any, Optional, List, Set
from datetime import datetime
from collections import defaultdict

logger = logging.getLogger(__name__)

# 文档解析依赖
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    import requests
    REQUESTS_AVAILABLE = True
except ImportError:
    REQUESTS_AVAILABLE = False


# ============================================================
# 多语言配置 - 参考 novel_config.py
# ============================================================
LANG_CONFIG = {
    "zh": {
        "name": "中文",
        "summary": "请为以下 {doc_type} 文档生成一段简洁的摘要（{summary_length} 字以内）：",
        "keywords": "请从以下文档内容中提取 5-8 个关键词（用逗号分隔）：",
        "relationships": "请分析以下文档之间的关系：",
        "overall": "请根据以下文档内容生成整体摘要：",
        "summary_label": "摘要",
        "keywords_label": "关键词",
        "relationship_label": "关系",
        "report_title": "文档学习报告",
        "generated_at": "生成时间",
        "total_docs": "共分析",
        "docs_label": "文档列表及摘要",
        "overall_label": "整体摘要",
        "relationship_graph": "文档关系图谱",
        "mermaid_label": "关系图（Mermaid）",
        "keywords_extract": "关键词提取",
        "appears": "出现次数",
    },
    "en": {
        "name": "English",
        "summary": "Please generate a concise summary for the following {doc_type} document (within {summary_length} words):",
        "keywords": "Please extract 5-8 keywords from the following document content (comma separated):",
        "relationships": "Please analyze the relationships between the following documents:",
        "overall": "Please generate an overall summary based on the following document content:",
        "summary_label": "Summary",
        "keywords_label": "Keywords",
        "relationship_label": "Relationship",
        "report_title": "Document Learning Report",
        "generated_at": "Generated At",
        "total_docs": "Total Documents",
        "docs_label": "Documents & Summaries",
        "overall_label": "Overall Summary",
        "relationship_graph": "Document Relationship Graph",
        "mermaid_label": "Relationship Diagram (Mermaid)",
        "keywords_extract": "Keyword Extraction",
        "appears": "Appears",
    },
    "ja": {
        "name": "日本語",
        "summary": "次の {doc_type} 文書の簡潔な要約を生成してください（{summary_length} 文字以内）：",
        "keywords": "以下の文書内容から 5-8 個のキーワードを抽出してください（カンマ区切り）：",
        "relationships": "以下の文書間の関係を分析してください：",
        "overall": "以下の文書内容に基づいて全体の要約を生成してください：",
        "summary_label": "要約",
        "keywords_label": "キーワード",
        "relationship_label": "関係",
        "report_title": "文書学習レポート",
        "generated_at": "生成時間",
        "total_docs": "合計文書数",
        "docs_label": "文書一覧と要約",
        "overall_label": "全体要約",
        "relationship_graph": "文書関係図",
        "mermaid_label": "関係図（Mermaid）",
        "keywords_extract": "キーワード抽出",
        "appears": "出現回数",
    },
    "ko": {
        "name": "한국어",
        "summary": "다음 {doc_type} 문서의 간결한 요약을 생성하세요（{summary_length} 자 이내）：",
        "keywords": "다음 문서 내용에서 5-8 개의 키워드를 추출하세요（쉼표로 구분）：",
        "relationships": "다음 문서 간의 관계를 분석하세요：",
        "overall": "다음 문서 내용을 기반으로 전체 요약을 생성하세요：",
        "summary_label": "요약",
        "keywords_label": "키워드",
        "relationship_label": "관계",
        "report_title": "문서 학습 보고서",
        "generated_at": "생성 시간",
        "total_docs": "총 문서 수",
        "docs_label": "문서 목록 및 요약",
        "overall_label": "전체 요약",
        "relationship_graph": "문서 관계도",
        "mermaid_label": "관계도（Mermaid）",
        "keywords_extract": "키워드 추출",
        "appears": "출현 횟수",
    },
    "es": {
        "name": "Español",
        "summary": "Por favor, genere un resumen conciso para el siguiente documento {doc_type} (dentro de {summary_length} palabras):",
        "keywords": "Por favor, extraiga 5-8 palabras clave del siguiente contenido del documento (separadas por comas):",
        "relationships": "Por favor, analice las relaciones entre los siguientes documentos:",
        "overall": "Por favor, genere un resumen general basado en el siguiente contenido del documento:",
        "summary_label": "Resumen",
        "keywords_label": "Palabras clave",
        "relationship_label": "Relación",
        "report_title": "Informe de Aprendizaje de Documentos",
        "generated_at": "Generado el",
        "total_docs": "Total de documentos",
        "docs_label": "Documentos y Resúmenes",
        "overall_label": "Resumen General",
        "relationship_graph": "Gráfico de Relaciones de Documentos",
        "mermaid_label": "Diagrama de Relaciones (Mermaid)",
        "keywords_extract": "Extracción de Palabras Clave",
        "appears": "Aparece",
    },
    "fr": {
        "name": "Français",
        "summary": "Veuillez générer un résumé concis pour le document {doc_type} suivant (dans la limite de {summary_length} mots) :",
        "keywords": "Veuillez extraire 5-8 mots-clés du contenu du document suivant (séparés par des virgules) :",
        "relationships": "Veuillez analyser les relations entre les documents suivants :",
        "overall": "Veuillez générer un résumé global basé sur le contenu du document suivant :",
        "summary_label": "Résumé",
        "keywords_label": "Mots-clés",
        "relationship_label": "Relation",
        "report_title": "Rapport d'Apprentissage de Documents",
        "generated_at": "Généré le",
        "total_docs": "Total de documents",
        "docs_label": "Documents et Résumés",
        "overall_label": "Résumé Global",
        "relationship_graph": "Graphe des Relations de Documents",
        "mermaid_label": "Diagramme des Relations (Mermaid)",
        "keywords_extract": "Extraction de Mots-clés",
        "appears": "Apparaît",
    },
    "de": {
        "name": "Deutsch",
        "summary": "Bitte erstellen Sie eine prägnante Zusammenfassung für das folgende {doc_type}-Dokument (innerhalb von {summary_length} Wörtern):",
        "keywords": "Bitte extrahieren Sie 5-8 Schlüsselwörter aus dem folgenden Dokumentinhalt (durch Kommas getrennt):",
        "relationships": "Bitte analysieren Sie die Beziehungen zwischen den folgenden Dokumenten:",
        "overall": "Bitte erstellen Sie eine Gesamtzusammenfassung basierend auf dem folgenden Dokumentinhalt:",
        "summary_label": "Zusammenfassung",
        "keywords_label": "Schlüsselwörter",
        "relationship_label": "Beziehung",
        "report_title": "Dokumenten-Lernbericht",
        "generated_at": "Erstellt am",
        "total_docs": "Dokumente insgesamt",
        "docs_label": "Dokumente und Zusammenfassungen",
        "overall_label": "Gesamtzusammenfassung",
        "relationship_graph": "Dokumenten-Beziehungsdiagramm",
        "mermaid_label": "Beziehungsdiagramm (Mermaid)",
        "keywords_extract": "Schlüsselwort-Extraktion",
        "appears": "Erscheint",
    },
}


class FastDocPointLearner:
    """
    文档快速学习助手
    """

    def __init__(self, config: Dict[str, Any] = None):
        self.config = config or {}
        self.name = "fast_doc_point_learner"
        self.version = "1.0.0"
        self._setup_logging()
        self._setup_config()

        self._check_dependencies()

        logger.info("FastDocPointLearner 初始化完成")

    def _setup_logging(self):
        log_level = self.config.get("log_level", "INFO")
        logging.basicConfig(
            level=getattr(logging, log_level.upper()),
            format="%(asctime)s - %(name)s - %(levelname)s - %(message)s"
        )

    def _setup_config(self):
        defaults = {
            "output_dir": "./skills/fast_doc_point_learner/output",
            "ollama_url": "http://localhost:11434",
            "model": "qwen2.5:7b",
            "max_chars": 5000,
            "summary_length": 300,
            "temperature": 0.3,
            "relationship_depth": 2,
            "extract_images": False,
            "extract_tables": True,
            "language": "zh",
        }
        for key, value in defaults.items():
            if key not in self.config:
                self.config[key] = value

        Path(self.config["output_dir"]).mkdir(parents=True, exist_ok=True)
        Path(self.config["output_dir"] + "/images").mkdir(parents=True, exist_ok=True)

    def _get_lang_config(self, lang: str) -> Dict:
        """获取语言配置"""
        if lang in LANG_CONFIG:
            return LANG_CONFIG[lang]
        logger.warning(f"不支持的语言: {lang}，使用中文")
        return LANG_CONFIG["zh"]

    def _check_dependencies(self):
        """检查文档解析依赖"""
        deps = {
            "DOCX": DOCX_AVAILABLE,
            "PDF": PDF_AVAILABLE,
            "PPTX": PPTX_AVAILABLE,
            "XLSX": XLSX_AVAILABLE,
        }
        for name, available in deps.items():
            if not available:
                logger.warning(f"{name} 解析不可用，请安装对应依赖")

    # ==================== 文档解析 ====================

    def _extract_text_from_docx(self, file_path: Path) -> str:
        """提取 Word 文档文本"""
        if not DOCX_AVAILABLE:
            return ""

        try:
            doc = Document(file_path)
            texts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    texts.append(para.text.strip())

            if self.config.get("extract_tables", True):
                for table in doc.tables:
                    for row in table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                        if row_text:
                            texts.append(row_text)

            return "\n".join(texts)
        except Exception as e:
            logger.warning(f"解析 DOCX 失败: {e}")
            return ""

    def _extract_text_from_pdf(self, file_path: Path) -> str:
        """提取 PDF 文档文本"""
        if not PDF_AVAILABLE:
            return ""

        try:
            texts = []
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        texts.append(text.strip())
            return "\n".join(texts)
        except Exception as e:
            logger.warning(f"解析 PDF 失败: {e}")
            return ""

    def _extract_text_from_pptx(self, file_path: Path) -> str:
        """提取 PPT 文档文本"""
        if not PPTX_AVAILABLE:
            return ""

        try:
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                    if hasattr(shape, "table") and self.config.get("extract_tables", True):
                        for row in shape.table.rows:
                            row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                            if row_text:
                                texts.append(row_text)
            return "\n".join(texts)
        except Exception as e:
            logger.warning(f"解析 PPTX 失败: {e}")
            return ""

    def _extract_text_from_xlsx(self, file_path: Path) -> str:
        """提取 Excel 文档文本"""
        if not XLSX_AVAILABLE:
            return ""

        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            texts = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                texts.append(f"=== 工作表: {sheet_name} ===")
                for row in sheet.iter_rows(values=True):
                    row_text = " | ".join([str(cell) for cell in row if cell is not None and str(cell).strip()])
                    if row_text:
                        texts.append(row_text)
            return "\n".join(texts)
        except Exception as e:
            logger.warning(f"解析 XLSX 失败: {e}")
            return ""

    def _extract_text_from_txt(self, file_path: Path) -> str:
        """提取 TXT 文档文本"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return content
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='gbk') as f:
                    content = f.read()
                return content
            except Exception as e:
                logger.warning(f"解析 TXT 失败: {e}")
                return ""

    def _extract_text_from_file(self, file_path: Path) -> tuple:
        """从文件提取文本和元数据"""
        suffix = file_path.suffix.lower()
        content = ""
        doc_type = ""

        if suffix == ".docx":
            content = self._extract_text_from_docx(file_path)
            doc_type = "Word"
        elif suffix == ".pdf":
            content = self._extract_text_from_pdf(file_path)
            doc_type = "PDF"
        elif suffix == ".pptx":
            content = self._extract_text_from_pptx(file_path)
            doc_type = "PPT"
        elif suffix == ".xlsx":
            content = self._extract_text_from_xlsx(file_path)
            doc_type = "Excel"
        elif suffix == ".txt":
            content = self._extract_text_from_txt(file_path)
            doc_type = "TXT"
        elif suffix == ".md":  # ✅ 新增 Markdown 支持
            content = self._extract_text_from_txt(file_path)
            doc_type = "Markdown"            
        else:
            return "", ""

        max_chars = self.config.get("max_chars", 5000)
        if len(content) > max_chars:
            content = content[:max_chars] + "\n\n... (内容已截断)"

        return content, doc_type

    # ==================== AI 摘要生成（多语言） ====================

    def _call_ollama(self, prompt: str, temperature: float = 0.3) -> str:
        """调用 Ollama API"""
        if not REQUESTS_AVAILABLE:
            return ""

        url = f"{self.config.get('ollama_url')}/api/generate"
        model = self.config.get("model", "qwen2.5:7b")

        payload = {
            "model": model,
            "prompt": prompt,
            "stream": False,
            "options": {
                "temperature": temperature,
                "num_predict": 1024,
            }
        }

        try:
            response = requests.post(url, json=payload, timeout=120)
            response.raise_for_status()
            data = response.json()
            return data.get("response", "").strip()
        except Exception as e:
            logger.error(f"Ollama 调用失败: {e}")
            return ""

    def _generate_summary(self, content: str, doc_name: str, doc_type: str) -> str:
        """生成文档摘要（多语言）"""
        if not content:
            return "无法提取内容"

        lang = self.config.get("language", "zh")
        lang_config = self._get_lang_config(lang)
        summary_length = self.config.get("summary_length", 300)

        prompt_template = lang_config.get("summary", "请为以下 {doc_type} 文档生成一段简洁的摘要（{summary_length} 字以内）：")
        prompt = prompt_template.format(doc_type=doc_type, summary_length=summary_length)
        prompt += f"\n\n文档名称：{doc_name}\n文档内容：\n{content[:3000]}\n\n请生成摘要："

        summary = self._call_ollama(prompt)
        if not summary:
            sentences = content.split("。")[:3]
            summary = "。".join(sentences) + "。" if sentences else "内容较短，无法生成摘要"

        return summary

    def _generate_keywords(self, content: str) -> List[str]:
        """提取关键词（多语言）"""
        if not content:
            return []

        lang = self.config.get("language", "zh")
        lang_config = self._get_lang_config(lang)
        prompt_template = lang_config.get("keywords", "请从以下文档内容中提取 5-8 个关键词（用逗号分隔）：")
        prompt = prompt_template + f"\n\n文档内容：\n{content[:2000]}\n\n关键词："

        keywords_text = self._call_ollama(prompt)
        if keywords_text:
            keywords = [k.strip() for k in keywords_text.split(",")]
            return keywords[:8]
        return []

    def _generate_relationships(self, doc_summaries: List[Dict]) -> List[Dict]:
        """分析文档之间的关系（多语言）"""
        if len(doc_summaries) < 2:
            return []

        lang = self.config.get("language", "zh")
        lang_config = self._get_lang_config(lang)
        relationships = []
        doc_info = "\n".join([f"- {d['name']}: {d['summary'][:200]}" for d in doc_summaries])

        prompt_template = lang_config.get("relationships", "请分析以下文档之间的关系：")
        prompt = f"{prompt_template}\n\n{doc_info}\n\n请分析文档之间的引用、依赖、补充等关系，输出 JSON 格式：\n{{'relationships': [{{'source': '文档A', 'target': '文档B', 'type': '引用|依赖|补充|扩展', 'description': '关系描述'}}]}}"

        result = self._call_ollama(prompt, temperature=0.1)

        try:
            json_match = re.search(r'\{.*\}', result, re.DOTALL)
            if json_match:
                data = json.loads(json_match.group().replace("'", '"'))
                return data.get("relationships", [])
        except:
            pass

        return relationships

    def _generate_overall_summary(self, doc_summaries: List[Dict]) -> str:
        """生成整体摘要（多语言）"""
        lang = self.config.get("language", "zh")
        lang_config = self._get_lang_config(lang)
        all_content = "\n".join([d["content"][:1000] for d in doc_summaries if d.get("content")])

        if not all_content:
            return f"{lang_config.get('total_docs', '共分析')} {len(doc_summaries)} 份文档。"

        prompt_template = lang_config.get("overall", "请根据以下文档内容生成整体摘要：")
        prompt = f"{prompt_template}\n\n{all_content[:3000]}\n\n整体摘要："

        summary = self._call_ollama(prompt)
        if summary:
            return summary
        else:
            lines = [f"{lang_config.get('total_docs', '共分析')} {len(doc_summaries)} 份文档。"]
            for doc in doc_summaries:
                lines.append(f"- {doc['name']}: {doc['summary'][:50]}...")
            return "\n".join(lines)

    # ==================== 文档收集 ====================

    def _collect_documents(self, doc_path: Path) -> List[Path]:
        """收集文档文件"""
        files = []
        supported_extensions = {".docx", ".pdf", ".pptx", ".xlsx", ".txt", ".md"}

        if doc_path.is_file():
            if doc_path.suffix.lower() in supported_extensions:
                files.append(doc_path)
        else:
            recursive = self.config.get("recursive", True)
            pattern = "**/*" if recursive else "*"
            for file_path in doc_path.glob(pattern):
                if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
                    if not file_path.name.startswith("~"):
                        files.append(file_path)

        return files

    # ==================== 生成报告（多语言） ====================

    def _generate_report(self, doc_summaries: List[Dict], relationships: List[Dict]) -> str:
        """生成 Markdown 报告（多语言）"""
        lang = self.config.get("language", "zh")
        lang_config = self._get_lang_config(lang)

        lines = []

        lines.append(f"# {lang_config.get('report_title', '文档学习报告')}")
        lines.append("")
        lines.append(f"> {lang_config.get('generated_at', '生成时间')}: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        lines.append("")
        lines.append(f"{lang_config.get('total_docs', '共分析')} **{len(doc_summaries)}** 份文档")
        lines.append("")

        # 文档列表
        lines.append(f"## {lang_config.get('docs_label', '文档列表及摘要')}")
        lines.append("")
        lines.append("| 文档 | 类型 | 摘要 | 关键词 |")
        lines.append("|------|------|------|--------|")
        for doc in doc_summaries:
            keywords = ", ".join(doc.get("keywords", [])[:3])
            summary_text = doc['summary'][:100]
            if len(doc['summary']) > 100:
                summary_text += "..."
            lines.append(f"| {doc['name']} | {doc['type']} | {summary_text} | {keywords} |")
        lines.append("")

        # 整体摘要
        lines.append(f"## {lang_config.get('overall_label', '整体摘要')}")
        lines.append("")
        overall_summary = self._generate_overall_summary(doc_summaries)
        lines.append(overall_summary)
        lines.append("")

        # 关系图谱
        if relationships:
            lines.append(f"## {lang_config.get('relationship_graph', '文档关系图谱')}")
            lines.append("")
            lines.append("```text")
            for rel in relationships:
                rel_type = rel.get('type', lang_config.get('relationship_label', '引用'))
                lines.append(f"{rel['source']} → {rel['target']} [{rel_type}]")
                lines.append(f"  {rel.get('description', '')}")
            lines.append("```")
            lines.append("")

            # Mermaid 图
            lines.append(f"### {lang_config.get('mermaid_label', '关系图（Mermaid）')}")
            lines.append("")
            lines.append("```mermaid")
            lines.append("graph TD")
            for rel in relationships:
                source = rel['source'].replace(" ", "_").replace(".", "_")
                target = rel['target'].replace(" ", "_").replace(".", "_")
                lines.append(f"    {source}[\"{rel['source']}\"] --> {target}[\"{rel['target']}\"]")
            lines.append("```")
            lines.append("")

        # 关键词提取
        all_keywords = {}
        for doc in doc_summaries:
            for kw in doc.get("keywords", []):
                all_keywords[kw] = all_keywords.get(kw, 0) + 1

        if all_keywords:
            lines.append(f"## {lang_config.get('keywords_extract', '关键词提取')}")
            lines.append("")
            lines.append(f"| {lang_config.get('keywords_label', '关键词')} | {lang_config.get('appears', '出现次数')} |")
            lines.append("|--------|----------|")
            sorted_kw = sorted(all_keywords.items(), key=lambda x: x[1], reverse=True)[:10]
            for kw, count in sorted_kw:
                lines.append(f"| {kw} | {count} |")
            lines.append("")

        lines.append("---")
        lines.append("")
        lines.append(f"*报告由 FastDocPointLearner 生成于 {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*")

        return "\n".join(lines)

    # ==================== 执行入口 ====================

    def execute(self, **kwargs) -> Dict[str, Any]:
        """执行文档学习"""
        logger.info(f"执行技能: {self.name} (v{self.version})")

        try:
            doc_path = kwargs.get("doc_path", "")
            if not doc_path:
                return {"status": "error", "error": "请提供 doc_path 参数"}

            # 语言配置
            lang = kwargs.get("language", self.config.get("language", "zh"))
            self.config["language"] = lang
            lang_config = self._get_lang_config(lang)
            logger.info(f"  语言: {lang_config.get('name', lang)} ({lang})")

            path = Path(doc_path)
            if not path.exists():
                return {"status": "error", "error": f"路径不存在: {doc_path}"}

            # 收集文档
            files = self._collect_documents(path)
            if not files:
                return {"status": "error", "error": "未找到支持的文档文件"}

            logger.info(f"找到 {len(files)} 个文档")

            # 解析每个文档
            doc_summaries = []
            for file_path in files:
                content, doc_type = self._extract_text_from_file(file_path)
                if not content:
                    continue

                summary = self._generate_summary(content, file_path.name, doc_type)
                keywords = self._generate_keywords(content)

                doc_summaries.append({
                    "name": file_path.name,
                    "path": str(file_path),
                    "type": doc_type,
                    "content": content,
                    "summary": summary,
                    "keywords": keywords,
                })

            if not doc_summaries:
                return {"status": "error", "error": "无法提取任何文档内容"}

            # 分析关系
            relationships = self._generate_relationships(doc_summaries)

            # 生成报告
            report = self._generate_report(doc_summaries, relationships)

            # 保存报告
            output_dir = Path(self.config["output_dir"])
            output_dir.mkdir(parents=True, exist_ok=True)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            report_file = output_dir / f"doc_learning_report_{timestamp}.md"

            with open(report_file, 'w', encoding='utf-8') as f:
                f.write(report)

            return {
                "status": "success",
                "result": {
                    "report_path": str(report_file),
                    "total_docs": len(doc_summaries),
                    "docs": doc_summaries,
                    "relationships": relationships,
                    "generated_at": datetime.now().isoformat(),
                },
                "metadata": {
                    "skill": self.name,
                    "version": self.version,
                }
            }

        except Exception as e:
            logger.error(f"执行失败: {e}")
            import traceback
            traceback.print_exc()
            return {
                "status": "error",
                "error": str(e),
                "skill": self.name,
            }

    def __repr__(self):
        return f"<FastDocPointLearner(name={self.name}, version={self.version})>"